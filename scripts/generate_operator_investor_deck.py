from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "investor-decks"
ASSETS = ROOT / "attached_assets"

BG = RGBColor(3, 21, 26)
PANEL = RGBColor(10, 42, 48)
PANEL_2 = RGBColor(13, 58, 65)
GOLD = RGBColor(248, 212, 119)
CYAN = RGBColor(34, 211, 238)
GREEN = RGBColor(52, 211, 153)
PINK = RGBColor(244, 114, 182)
VIOLET = RGBColor(167, 139, 250)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(185, 202, 207)
LINE = RGBColor(31, 82, 90)

FONT = "Aptos"
ADDRESS = "Ruko Oceanic Bliss, Batam Centre, Batam, Indonesia"
MAP_URL = (
    "https://www.google.com/maps/search/?api=1&query=Ruko%20Oceanic%20Bliss%2C%20Jl.%20Pasir%20Putih%20Harbourfront%20-%20Batam%20Centre.49%20blok%20A.%2051%2C%20Sadai%2C%20Bengkong%2C%20Batam%20City%2C%20Riau%20Islands%2029444"
)
DISCLAIMER = (
    "ROI projections are estimates, not guaranteed returns. Figures require review by qualified "
    "legal, tax, and accounting professionals before operator contracts are signed."
)


BUSINESSES = [
    {
        "code": "1FA",
        "name": "Kids Fun House",
        "units": 1,
        "sqm": 75,
        "setup_idr": "Rp875M",
        "setup_sgd": "S$63K",
        "revenue": "Rp140M",
        "opex": "Rp95M",
        "net": "Rp45M",
        "operator_profit": "Rp378M",
        "roi": "43%",
        "why": "Family and tourist daytime traffic, prize games, birthday visits, and safe kid-friendly activities.",
        "accent": PINK,
    },
    {
        "code": "1FB",
        "name": "KTV Lounge",
        "units": 2,
        "sqm": 150,
        "setup_idr": "Rp1.75B",
        "setup_sgd": "S$126K",
        "revenue": "Rp300M",
        "opex": "Rp205M",
        "net": "Rp95M",
        "operator_profit": "Rp798M",
        "roi": "46%",
        "why": "Open singing, group seating, singing competitions, and social video moments for repeat traffic.",
        "accent": CYAN,
    },
    {
        "code": "2FA",
        "name": "Beauty Salon",
        "units": 1,
        "sqm": 75,
        "setup_idr": "Rp775M",
        "setup_sgd": "S$56K",
        "revenue": "Rp135M",
        "opex": "Rp85M",
        "net": "Rp50M",
        "operator_profit": "Rp420M",
        "roi": "54%",
        "why": "Repeat appointments, group makeovers, event prep, and cross-traffic from KTV guests.",
        "accent": PINK,
    },
    {
        "code": "2FB",
        "name": "Private KTV Rooms",
        "units": 2,
        "sqm": 150,
        "setup_idr": "Rp1.9B",
        "setup_sgd": "S$137K",
        "revenue": "Rp350M",
        "opex": "Rp235M",
        "net": "Rp115M",
        "operator_profit": "Rp966M",
        "roi": "51%",
        "why": "Four bookable rooms for birthdays, friend groups, customer hosting, and private celebrations.",
        "accent": CYAN,
    },
    {
        "code": "3FA",
        "name": "Beauty Salon",
        "units": 1,
        "sqm": 75,
        "setup_idr": "Rp775M",
        "setup_sgd": "S$56K",
        "revenue": "Rp145M",
        "opex": "Rp90M",
        "net": "Rp55M",
        "operator_profit": "Rp462M",
        "roi": "60%",
        "why": "Beauty service beside VIP KTV supports premium guests, event prep, and recurring spend.",
        "accent": PINK,
    },
    {
        "code": "3FB",
        "name": "VIP KTV Rooms",
        "units": 2,
        "sqm": 150,
        "setup_idr": "Rp2.15B",
        "setup_sgd": "S$155K",
        "revenue": "Rp400M",
        "opex": "Rp270M",
        "net": "Rp130M",
        "operator_profit": "Rp1.092B",
        "roi": "51%",
        "why": "Two premium rooms for VIP groups, business hosting, birthday events, and higher table spend.",
        "accent": GOLD,
    },
    {
        "code": "4F",
        "name": "Pet Cafe",
        "units": 3,
        "sqm": 225,
        "setup_idr": "Rp2.5B",
        "setup_sgd": "S$180K",
        "revenue": "Rp450M",
        "opex": "Rp330M",
        "net": "Rp120M",
        "operator_profit": "Rp1.008B",
        "roi": "40%",
        "why": "Food, drinks, pet community, blindbox tie-ins, photo content, and family-friendly visits.",
        "accent": GREEN,
    },
    {
        "code": "5F",
        "name": "Live House",
        "units": 3,
        "sqm": 225,
        "setup_idr": "Rp3.6B",
        "setup_sgd": "S$259K",
        "revenue": "Rp750M",
        "opex": "Rp550M",
        "net": "Rp200M",
        "operator_profit": "Rp1.68B",
        "roi": "47%",
        "why": "Live band, performances, dance floor, sea-view events, premium tables, and nightlife traffic.",
        "accent": VIOLET,
    },
]


def tx(slide, text, x, y, w, h, size=24, color=WHITE, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def rect(slide, x, y, w, h, fill=PANEL, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line or LINE
    shp.line.transparency = 25
    return shp


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    rect(slide, 10.55, -0.2, 2.4, 7.9, RGBColor(8, 55, 61), RGBColor(8, 55, 61), False)
    for i, color in enumerate([CYAN, GOLD, GREEN]):
        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2 + i * 0.55), Inches(0.45 + i), Inches(2.3), Inches(2.3))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.fill.transparency = 83
        shp.line.transparency = 100


def image(slide, filename, x, y, w, h=None):
    path = ASSETS / filename
    if not path.exists():
        return
    if h:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))


def footer(slide):
    tx(slide, "Reborn Wave Group | Operator-Investor Deck | IDR main + SGD reference", 0.55, 7.05, 7.8, 0.2, 7, MUTED)
    tx(slide, "ROI estimates only, not guaranteed returns.", 8.2, 6.95, 4.6, 0.25, 6, MUTED, align=PP_ALIGN.RIGHT)


def title(slide, heading, sub=None):
    tx(slide, "REBORN WAVE GROUP | OPERATOR-INVESTOR DECK", 0.55, 0.25, 4.2, 0.28, 9, MUTED, True)
    tx(slide, heading, 0.55, 0.62, 9.3, 0.72, 30, WHITE, True)
    if sub:
        tx(slide, sub, 0.58, 1.28, 8.8, 0.5, 12, MUTED)


def card(slide, x, y, w, h, heading, body, accent=CYAN, body_size=9):
    rect(slide, x, y, w, h, PANEL)
    tx(slide, heading, x + 0.18, y + 0.15, w - 0.36, 0.32, 13, accent, True)
    tx(slide, body, x + 0.18, y + 0.56, w - 0.36, h - 0.68, body_size, MUTED)


def metric(slide, x, y, w, label, value, accent=GOLD):
    rect(slide, x, y, w, 0.78, RGBColor(7, 39, 45), accent)
    tx(slide, value, x + 0.12, y + 0.12, w - 0.24, 0.28, 16, accent, True, PP_ALIGN.CENTER)
    tx(slide, label, x + 0.12, y + 0.45, w - 0.24, 0.2, 7, MUTED, False, PP_ALIGN.CENTER)


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    image(slide, "rwg-logo.png", 0.65, 0.42, 1.18)
    tx(slide, "Reborn Wave Group", 0.72, 1.72, 6.2, 0.7, 39, WHITE, True)
    tx(slide, "5-Level Batam Activity Club", 0.76, 2.58, 5.9, 0.38, 18, GOLD, True)
    tx(slide, "Operator-Investor Opportunity: 2 Years $0 Rent", 0.76, 3.25, 5.6, 0.55, 16, MUTED)
    tx(slide, "Business owners fund renovation and operations, then keep 70% of monthly net profit.", 0.76, 4.08, 5.7, 0.65, 13, WHITE)
    rect(slide, 7.05, 0.9, 5.15, 5.6, PANEL_2)
    image(slide, "doluruu-blindbox-box.jpeg", 7.38, 1.14, 4.45, 2.55)
    image(slide, "doluruu-female-transparent.png", 8.18, 3.38, 1.45)
    image(slide, "Doluruu Boy_1749664545355.png", 9.62, 3.35, 1.32)
    tx(slide, "75 sqm per unit | IDR main + SGD reference", 7.42, 5.5, 4.45, 0.35, 14, GOLD, True, PP_ALIGN.CENTER)
    footer(slide)


def opportunity(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "The Opportunity", "Tourism, family, nightlife, beauty, KTV, pet cafe, and events in one building.")
    items = [
        ("One Destination", "A five-level Batam club where tourists, families, groups, and nightlife guests can stay longer and spend across multiple activities."),
        ("Partner-Led Growth", "Operators bring capital and operating expertise while Reborn Wave provides space, brand, traffic, website, social media, and shared customer flow."),
        ("Cross-Traffic", "KTV, kids play, beauty, pet cafe, and live events feed each other through packages, birthdays, competitions, and social content."),
    ]
    for i, item in enumerate(items):
        card(slide, 0.7 + i * 4.05, 1.95, 3.55, 2.6, item[0], item[1], [CYAN, GOLD, GREEN][i], 10)
    footer(slide)


def operator_deal(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "Operator Deal", "A lower-rent-risk path for business owners to operate inside a ready destination.")
    deal = [
        ("$0 rent for 24 months", "Operator gets two years rent-free from opening or handover date."),
        ("Operator funds the business", "Renovation, equipment, staff, stock, licenses, utilities, and operations are covered by the operator."),
        ("70/30 monthly profit split", "Monthly net operating profit is split 70% to operator and 30% to Reborn Wave Group."),
        ("Transparent reporting", "Sales run through shared POS or audited monthly report; cash sales must be recorded daily."),
    ]
    for i, item in enumerate(deal):
        card(slide, 0.7 + (i % 2) * 6.1, 1.72 + (i // 2) * 1.85, 5.65, 1.42, item[0], item[1], [GOLD, CYAN, GREEN, PINK][i], 10)
    footer(slide)


def floor_breakdown(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "Floor + Unit Breakdown", "Updated area assumption: 75 sqm per unit.")
    for i, business in enumerate(BUSINESSES):
        y = 1.62 + i * 0.66
        rect(slide, 0.7, y, 11.75, 0.5, RGBColor(7, 39, 45), business["accent"])
        tx(slide, business["code"], 0.9, y + 0.1, 0.65, 0.18, 12, business["accent"], True)
        tx(slide, business["name"], 1.55, y + 0.08, 2.65, 0.22, 11, WHITE, True)
        tx(slide, f'{business["units"]} unit(s)', 4.35, y + 0.09, 1.25, 0.2, 9, MUTED)
        tx(slide, f'{business["sqm"]} sqm', 5.72, y + 0.09, 1.1, 0.2, 9, GOLD, True)
        tx(slide, business["why"], 6.9, y + 0.09, 5.35, 0.22, 8, MUTED)
    footer(slide)


def roi_table(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "Budget + ROI Projection", "Balanced base-case projections. Operator annual ROI is based on 70% profit share.")
    headers = ["Business", "Setup IDR", "Setup SGD", "Monthly Net", "Operator Annual", "ROI"]
    widths = [2.6, 1.65, 1.35, 1.75, 1.9, 0.9]
    x0, y0 = 0.55, 1.55
    x = x0
    for header, width in zip(headers, widths):
        rect(slide, x, y0, width, 0.42, RGBColor(8, 61, 68), CYAN, False)
        tx(slide, header, x + 0.05, y0 + 0.1, width - 0.1, 0.13, 7, WHITE, True, PP_ALIGN.CENTER)
        x += width
    for r, b in enumerate(BUSINESSES):
        y = y0 + 0.43 + r * 0.52
        vals = [f'{b["code"]} {b["name"]}', b["setup_idr"], b["setup_sgd"], b["net"], b["operator_profit"], b["roi"]]
        x = x0
        for c, (val, width) in enumerate(zip(vals, widths)):
            fill = RGBColor(7, 39, 45) if r % 2 == 0 else RGBColor(9, 48, 54)
            rect(slide, x, y, width, 0.5, fill, LINE, False)
            color = b["accent"] if c == 5 else WHITE if c == 0 else MUTED
            tx(slide, val, x + 0.05, y + 0.12, width - 0.1, 0.16, 7, color, c in [0, 5], PP_ALIGN.CENTER if c > 0 else None)
            x += width
    tx(slide, "Exchange assumption: SGD 1 approx. IDR 13,900.", 0.65, 6.45, 6.3, 0.25, 8, MUTED)
    footer(slide)


def business_slide(prs, b):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, f'{b["code"]} {b["name"]}', f'{b["units"]} unit(s) | {b["sqm"]} sqm | 70% operator / 30% Reborn Wave')
    metric(slide, 0.72, 1.72, 1.6, "Setup IDR", b["setup_idr"], b["accent"])
    metric(slide, 2.55, 1.72, 1.45, "Setup SGD", b["setup_sgd"], GOLD)
    metric(slide, 4.25, 1.72, 1.55, "Monthly net", b["net"], CYAN)
    metric(slide, 6.05, 1.72, 1.65, "Annual operator", b["operator_profit"], GREEN)
    metric(slide, 7.95, 1.72, 1.25, "Annual ROI", b["roi"], b["accent"])
    card(slide, 0.72, 3.0, 4.0, 2.05, "Monthly assumptions", f'Revenue: {b["revenue"]}\nOperating costs: {b["opex"]}\nNet profit before split: {b["net"]}', CYAN, 12)
    card(slide, 5.0, 3.0, 3.65, 2.05, "Why it attracts customers", b["why"], b["accent"], 11)
    card(slide, 8.95, 3.0, 3.0, 2.05, "Operator upside", f'Operator receives 70% of net monthly profit.\nAnnual ROI projection: {b["roi"]}.', GOLD, 11)
    footer(slide)


def reborn_projection(prs):
    monthly = "Rp255M/month"
    annual = "Rp3.06B/year"
    sgd = "S$220K/year"
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "Reborn Wave Revenue Projection", "Reborn receives 30% of each operator's monthly net profit.")
    metric(slide, 0.85, 1.8, 2.8, "Monthly Reborn profit share", monthly, GOLD)
    metric(slide, 4.0, 1.8, 2.8, "Annual Reborn profit share", annual, CYAN)
    metric(slide, 7.15, 1.8, 2.4, "SGD reference", sgd, GREEN)
    card(slide, 0.85, 3.25, 5.2, 1.7, "Why the model works", "Operators carry setup and operating costs. Reborn monetizes location, brand, traffic, systems, and shared customer flow without charging rent during the first two years.", CYAN, 11)
    card(slide, 6.45, 3.25, 4.75, 1.7, "Portfolio logic", "Multiple verticals reduce reliance on one business type: family traffic, beauty appointments, KTV bookings, pet cafe visits, and live events.", GOLD, 11)
    footer(slide)


def marketing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "Marketing + Traffic Strategy", "Create repeat reasons to visit, post, book, and bring groups.")
    items = [
        ("Singing competitions", "Weekly competition content, finalists, voting, prize moments, and live audience energy."),
        ("Tourist + family packages", "Bundled routes for kids fun house, KTV, pet cafe, beauty, and live house nights."),
        ("Birthday + private events", "Private KTV rooms, VIP hosting, live band packages, and table reservations."),
        ("Performance calendar", "Live bands, dancers, creators, guest singers, and sea-view theme nights."),
        ("TikTok", "@reborn.wave.group"),
        ("Instagram", "@rebornwavegroup"),
    ]
    for i, item in enumerate(items):
        card(slide, 0.7 + (i % 3) * 4.05, 1.65 + (i // 3) * 1.75, 3.6, 1.28, item[0], item[1], [CYAN, GOLD, GREEN, PINK, CYAN, VIOLET][i], 10)
    footer(slide)


def risk_controls(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "Risk Controls", "Protect operators, Reborn Wave, and the brand experience.")
    controls = [
        ("Shared POS / reporting", "Monthly net profit can be verified from sales, costs, and cash records."),
        ("Minimum operating hours", "Operators must keep the destination active and predictable for customers."),
        ("Brand + safety rules", "Design, hygiene, service, customer safety, and noise rules protect the building."),
        ("Monthly performance review", "Slow operators receive action plans before replacement rights are triggered."),
        ("Replacement right", "Inactive or non-performing operators can be replaced under contract terms."),
    ]
    for i, item in enumerate(controls):
        card(slide, 0.85, 1.55 + i * 0.92, 10.9, 0.72, item[0], item[1], [CYAN, GOLD, GREEN, PINK, VIOLET][i], 9)
    footer(slide)


def next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "Next Steps", "Move from operator interest to signed launch partners.")
    steps = [
        ("1. Confirm operators by business type", "Match each operator to the right floor zone and budget requirement."),
        ("2. Sign LOI", "Lock intent, contribution scope, timeline, reporting rules, and 70/30 split principles."),
        ("3. Finalize renovation scope", "Confirm design, cost, contractor timeline, equipment list, and opening budget."),
        ("4. Open pilot marketing campaign", "Use TikTok, Instagram, WhatsApp, creators, hotel partners, and event previews."),
        ("5. Launch by floor / zone", "Open staged areas as operators finish buildout, training, and systems setup."),
    ]
    for i, item in enumerate(steps):
        card(slide, 0.85, 1.45 + i * 0.92, 10.9, 0.72, item[0], item[1], [CYAN, GOLD, GREEN, PINK, VIOLET][i], 9)
    tx(slide, DISCLAIMER, 0.85, 6.18, 10.8, 0.45, 8, MUTED)
    footer(slide)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    cover(prs)
    opportunity(prs)
    operator_deal(prs)
    floor_breakdown(prs)
    roi_table(prs)
    for business in BUSINESSES:
        business_slide(prs, business)
    reborn_projection(prs)
    marketing(prs)
    risk_controls(prs)
    next_steps(prs)
    OUT.mkdir(exist_ok=True)
    output = OUT / "Reborn-Wave-Operator-Investor-Deck-EN.pptx"
    prs.save(output)
    return output


if __name__ == "__main__":
    print(build())
