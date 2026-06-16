from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "investor-decks"
ASSETS = ROOT / "attached_assets"

BG = RGBColor(3, 21, 26)
PANEL = RGBColor(10, 42, 48)
PANEL_2 = RGBColor(13, 58, 65)
GOLD = RGBColor(248, 212, 119)
CYAN = RGBColor(34, 211, 238)
GREEN = RGBColor(52, 211, 153)
PINK = RGBColor(244, 114, 182)
VIOLET = RGBColor(167, 139, 250)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(185, 202, 207)
LINE = RGBColor(31, 82, 90)
FONT = "Microsoft YaHei"

DISCLAIMER = (
    "本简报用于说明计划中的投资系统。回报、Token 价值、推荐奖励与利润分配均为预测，"
    "不构成保证。正式推出前需完成法律、税务、会计与牌照合规审核。"
)

PACKAGES = [
    ("$500", "$500 消费额度", "1 股", "10% 推荐奖励", "-"),
    ("$1,000", "$1,000 消费额度", "2 股", "10% 推荐奖励", "2天1夜公司参观"),
    ("$3,000", "$3,000 消费额度", "6 股", "10% 推荐奖励", "3天2夜参观 + 娱乐"),
    ("$5,000", "$7,000 消费额度", "10 股", "10% 推荐奖励", "3天2夜参观 + 娱乐"),
    ("$10,000", "$15,000 消费额度", "20 股", "10% 推荐奖励", "3天2夜参观 + 娱乐"),
]

PROJECTS = [
    ("KTV Lounge", "唱歌、包厢、接待、比赛活动"),
    ("Beauty", "美容美发、妆造、活动前造型"),
    ("Game House", "儿童、家庭游戏、兑换玩法"),
    ("Live House", "海景乐队、舞池、现场表演"),
    ("Pet Cafe", "餐饮、宠物、盲盒内容与社群"),
]


def tx(slide, text, x, y, w, h, size=18, color=WHITE, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def rect(slide, x, y, w, h, fill=PANEL, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line or LINE
    shp.line.transparency = 22
    return shp


def image(slide, filename, x, y, w, h=None):
    path = ASSETS / filename
    if not path.exists():
        return
    if h:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    rect(slide, 10.55, -0.2, 2.4, 7.9, RGBColor(8, 55, 61), RGBColor(8, 55, 61), False)
    for i, color in enumerate([CYAN, GOLD, GREEN]):
        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2 + i * 0.55), Inches(0.45 + i), Inches(2.3), Inches(2.3))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.fill.transparency = 83
        shp.line.transparency = 100


def footer(slide):
    tx(slide, "Reborn Wave Group | 投资配套 + Token 系统", 0.55, 7.05, 7.1, 0.2, 7, MUTED)
    tx(slide, "预测回报，不保证收益。", 8.3, 6.95, 4.3, 0.24, 6, MUTED, align=PP_ALIGN.RIGHT)


def title(slide, heading, sub=None):
    tx(slide, "REBORN WAVE GROUP | 普通投资人指南", 0.55, 0.25, 4.8, 0.28, 9, MUTED, True)
    tx(slide, heading, 0.55, 0.62, 9.4, 0.72, 30, WHITE, True)
    if sub:
        tx(slide, sub, 0.58, 1.28, 9.2, 0.5, 12, MUTED)


def card(slide, x, y, w, h, heading, body, accent=CYAN, body_size=9):
    rect(slide, x, y, w, h, PANEL)
    tx(slide, heading, x + 0.18, y + 0.14, w - 0.36, 0.32, 13, accent, True)
    tx(slide, body, x + 0.18, y + 0.54, w - 0.36, h - 0.62, body_size, MUTED)


def metric(slide, x, y, w, label, value, accent=GOLD):
    rect(slide, x, y, w, 0.78, RGBColor(7, 39, 45), accent)
    tx(slide, value, x + 0.12, y + 0.12, w - 0.24, 0.28, 16, accent, True, PP_ALIGN.CENTER)
    tx(slide, label, x + 0.12, y + 0.45, w - 0.24, 0.2, 7, MUTED, False, PP_ALIGN.CENTER)


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    image(slide, "rwg-logo.png", 0.65, 0.42, 1.18)
    tx(slide, "Reborn Wave Group", 0.72, 1.55, 6.2, 0.7, 38, WHITE, True)
    tx(slide, "$500 - $10,000 投资配套指南", 0.76, 2.44, 6.5, 0.45, 18, GOLD, True)
    tx(slide, "选择你看好的业务，获得每日项目 Token，参与每月业务利润池，并随着更多投资人加入推动 Token 池价值成长。", 0.76, 3.14, 5.95, 1.05, 13, WHITE)
    rect(slide, 7.05, 0.9, 5.15, 5.6, PANEL_2)
    image(slide, "image_1764558901062.png", 7.38, 1.14, 4.45, 2.42)
    image(slide, "doluruu-blindbox-box.jpeg", 7.55, 3.88, 2.1, 1.22)
    image(slide, "doluruu-female-transparent.png", 10.0, 3.6, 1.12)
    tx(slide, "投资入口: rebornwave.group/investor", 7.42, 5.55, 4.45, 0.35, 14, GOLD, True, PP_ALIGN.CENTER)
    footer(slide)


def package_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "投资配套", "每个配套包含俱乐部消费额度、项目股份、推荐奖励，以及部分配套的参观娱乐礼遇。")
    headers = ["金额", "消费额度", "股份", "推荐", "额外礼遇"]
    widths = [1.35, 2.0, 1.35, 1.65, 4.15]
    x0, y0 = 0.82, 1.75
    x = x0
    for h, w in zip(headers, widths):
        rect(slide, x, y0, w, 0.5, RGBColor(8, 61, 68), CYAN, False)
        tx(slide, h, x + 0.05, y0 + 0.13, w - 0.1, 0.16, 8, WHITE, True, PP_ALIGN.CENTER)
        x += w
    for i, row in enumerate(PACKAGES):
        y = y0 + 0.55 + i * 0.62
        x = x0
        for c, (val, w) in enumerate(zip(row, widths)):
            fill = RGBColor(7, 39, 45) if i % 2 == 0 else RGBColor(9, 48, 54)
            rect(slide, x, y, w, 0.58, fill, LINE, False)
            color = GOLD if c in [1, 2] else WHITE if c == 0 else MUTED
            tx(slide, val, x + 0.06, y + 0.15, w - 0.12, 0.16, 8, color, c in [0, 1, 2], PP_ALIGN.CENTER if c < 3 else None)
            x += w
    card(slide, 0.82, 5.55, 10.65, 0.72, "共享池潜在收益", "每股以 $500 参与额计算，有效期两年。公司 40% 共享池分配给参与投资人，并有机会获得 5-10 倍潜在回报；此为预测，不保证收益。", GOLD, 8)
    footer(slide)


def choose_business(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "选择要投资的业务", "投资人用分配单位投票选择业务。选择的业务决定每日获得哪一种项目 Token。")
    for i, (name, desc) in enumerate(PROJECTS):
        card(slide, 0.72 + (i % 3) * 4.05, 1.75 + (i // 3) * 1.65, 3.55, 1.18, name, desc, [CYAN, PINK, GOLD, VIOLET, GREEN][i], 9)
    card(slide, 0.72, 5.18, 11.25, 0.86, "分配规则", "投资配套的单位必须完整分配。例子: $10,000 Anchor Pack 有 20 单位，投资人可决定每个业务放多少单位。", GOLD, 9)
    footer(slide)


def token_engine(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "每日 Token + Token 价值成长", "Token 用来奖励参与，并把投资人与所选择的业务连接起来。")
    metric(slide, 0.85, 1.88, 2.4, "每日规则", "1 单位 = 1 Token/天", CYAN)
    metric(slide, 3.65, 1.88, 2.4, "Token 价格", "资金池 / 活跃 Token", GOLD)
    metric(slide, 6.45, 1.88, 2.4, "配套资金池", "价值成长", GREEN)
    metric(slide, 9.25, 1.88, 2.4, "出售选项", "Token -> 现金", PINK)
    card(slide, 0.85, 3.38, 3.35, 1.35, "每日领取", "投资人在后台点击领取每日项目 Token。Token 数量按活跃分配单位计算。", CYAN, 9)
    card(slide, 4.55, 3.38, 3.35, 1.35, "价值可成长", "当更多人加入时，部分配套活动进入全局 Token 池，支持 Token 价格。", GOLD, 9)
    card(slide, 8.25, 3.38, 3.35, 1.35, "出售 Token", "当资金池有流动性时，投资人可把符合条件的 Token 卖回现金钱包。", GREEN, 9)
    card(slide, 0.85, 5.25, 10.75, 0.72, "简单例子", "$5,000 Founder Pack = 10 单位。若 4 单位放 KTV、6 单位放 Live House，每日可领 4 KTV Token + 6 Live House Token。", VIOLET, 9)
    footer(slide)


def profit_pool(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "每月业务利润池", "业务利润模型奖励投票支持该业务的投资人。")
    metric(slide, 0.95, 1.86, 2.55, "每月核实利润", "业务报告", CYAN)
    metric(slide, 3.9, 1.86, 2.55, "投资人池", "40%", GOLD)
    metric(slide, 6.85, 1.86, 2.55, "分配对象", "投票者", GREEN)
    metric(slide, 9.8, 1.86, 1.7, "依据", "单位", PINK)
    card(slide, 0.85, 3.35, 3.35, 1.4, "谁可获得", "把单位分配到该业务的投资人，会被视为该业务的投票者。", CYAN, 9)
    card(slide, 4.55, 3.35, 3.35, 1.4, "如何分配", "每月净利润的 40% 按该业务活跃单位/投票比例分配。", GOLD, 9)
    card(slide, 8.25, 3.35, 3.35, 1.4, "后台审核", "审核销售、成本、净利润、投资人单位总数与 payout 批准。", GREEN, 9)
    card(slide, 0.85, 5.25, 10.75, 0.72, "例子", "若 Pet Cafe 每月核实净利润为 $10,000，则 $4,000 进入 Pet Cafe 投票者利润池，并按 Pet Cafe 单位分配。", VIOLET, 9)
    footer(slide)


def referral_system(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "推荐系统", "投资人可邀请别人加入，并在被推荐人购买配套时获得奖励。")
    metric(slide, 0.95, 1.86, 2.55, "一级奖励", "10%", GOLD)
    metric(slide, 3.9, 1.86, 2.55, "二级奖励", "5%", CYAN)
    metric(slide, 6.85, 1.86, 2.55, "奖励钱包", "现金 + 消费", GREEN)
    metric(slide, 9.8, 1.86, 1.7, "推荐码", "INV-xxxx", PINK)
    card(slide, 0.85, 3.35, 3.35, 1.35, "一级", "直接邀请人可在被邀请投资人购买配套后获得推荐奖励。", GOLD, 9)
    card(slide, 4.55, 3.35, 3.35, 1.35, "二级", "符合条件时，第二层邀请人可获得较小比例的奖励。", CYAN, 9)
    card(slide, 8.25, 3.35, 3.35, 1.35, "支持资金池", "无人领取的推荐奖励可回流 Token 池，支持 Token 价值。", GREEN, 9)
    card(slide, 0.85, 5.2, 10.75, 0.78, "奖励分配", "系统会把推荐奖励记录到现金钱包与消费钱包，让投资人一部分可提现，一部分可在俱乐部使用。", VIOLET, 9)
    footer(slide)


def topup_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "如何充值", "投资人后台支持 USDT BEP20 充值，并由管理员确认入账。")
    steps = [
        ("1. 注册 / 登录", "进入 rebornwave.group/investor/login 打开投资人后台。"),
        ("2. 选择配套", "选择 $500、$1,000、$3,000、$5,000 或 $10,000 配套。"),
        ("3. 连接钱包", "连接 BEP20 钱包，或在管理员允许时使用人工付款凭证。"),
        ("4. 转 USDT", "把准确金额转到 Reborn Wave 管理员钱包。"),
        ("5. 提交 TX Hash", "在后台粘贴交易哈希等待确认。"),
        ("6. 现金钱包入账", "确认后，现金余额入账，可用于购买投资配套。"),
    ]
    for i, (h, b) in enumerate(steps):
        card(slide, 0.72 + (i % 3) * 4.05, 1.75 + (i // 3) * 1.72, 3.55, 1.22, h, b, [CYAN, GOLD, GREEN, PINK, VIOLET, CYAN][i], 8)
    card(slide, 0.72, 5.48, 11.25, 0.78, "安全提醒", "只使用投资人后台显示或 Reborn Wave 官方管理员确认的钱包地址。", GOLD, 9)
    footer(slide)


def purchase_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "如何购买配套 + 分配单位", "充值后，投资人购买配套并选择项目分配。")
    card(slide, 0.72, 1.78, 3.55, 1.32, "选择配套", "选择金额后，系统会按每 $500 自动计算单位数。", CYAN, 9)
    card(slide, 4.55, 1.78, 3.55, 1.32, "分配单位", "购买确认前，选择每个业务项目要分配多少单位。", GOLD, 9)
    card(slide, 8.38, 1.78, 3.55, 1.32, "确认购买", "现金钱包支付配套，后台记录购买与单位分配。", GREEN, 9)
    card(slide, 0.72, 3.68, 5.45, 1.36, "例子 A", "$1,000 Builder Pack = 2 单位。投资人选择 1 KTV 单位 + 1 Pet Cafe 单位。", PINK, 10)
    card(slide, 6.48, 3.68, 5.45, 1.36, "例子 B", "$10,000 Anchor Pack = 20 单位。投资人选择 8 Live House、6 KTV、4 Pet Cafe、2 Beauty。", VIOLET, 10)
    card(slide, 0.72, 5.55, 11.25, 0.72, "购买后", "投资人可每日领取项目 Token、查看推荐、出售 Token、申请提现，并用 QR 钱包在俱乐部消费。", CYAN, 9)
    footer(slide)


def dashboard_usage(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "如何使用投资人后台", "后台用于管理充值、配套、Token、消费、推荐与提现。")
    items = [
        ("Packages", "选择配套并把单位分配到指定业务。"),
        ("Top-ups", "使用 USDT 交易哈希或管理员确认付款来入账。"),
        ("Tokens", "每日领取项目 Token，并查看活跃/已售 Token。"),
        ("Sell tokens", "按当前资金池价格，把符合条件 Token 卖到现金钱包。"),
        ("QR spending", "用消费钱包在俱乐部内扫码消费。"),
        ("Withdrawals", "申请把现金钱包提现到批准的钱包地址。"),
    ]
    for i, (h, b) in enumerate(items):
        card(slide, 0.72 + (i % 3) * 4.05, 1.75 + (i // 3) * 1.72, 3.55, 1.22, h, b, [CYAN, GOLD, GREEN, PINK, VIOLET, CYAN][i], 9)
    card(slide, 0.72, 5.5, 11.25, 0.72, "后台链接", "投资首页: rebornwave.group/investor   |   登录/注册: rebornwave.group/investor/login   |   后台: rebornwave.group/investor/dashboard", GOLD, 9)
    footer(slide)


def risk_controls(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "规则、报告与风险控制", "清晰记录能保护投资人、运营方与 Reborn Wave Group。")
    card(slide, 0.72, 1.78, 3.55, 1.35, "每月报告", "利润分配应以核实后的收入、成本、净利润与批准后的计算为准。", CYAN, 9)
    card(slide, 4.55, 1.78, 3.55, 1.35, "Token 池", "Token 价值取决于资金池余额与活跃 Token 供应，价格可能上涨也可能下降。", GOLD, 9)
    card(slide, 8.38, 1.78, 3.55, 1.35, "管理员批准", "充值、提现、配套设置与业务报告都需要后台管理控制。", GREEN, 9)
    card(slide, 0.72, 3.68, 5.45, 1.36, "不保证收益", "更多投资人加入可推动资金池价值，但收益不保证，取决于真实系统活动。", PINK, 9)
    card(slide, 6.48, 3.68, 5.45, 1.36, "专业审核", "正式公开前，合约、证券规则、税务、会计、牌照与投资风险披露必须审核。", VIOLET, 9)
    tx(slide, DISCLAIMER, 0.85, 5.72, 10.8, 0.48, 8, MUTED)
    footer(slide)


def next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "投资人流程总结", "从注册账户到购买配套与每日领取奖励。")
    steps = [
        ("1. 创建账户", "在 rebornwave.group/investor/login 注册。"),
        ("2. 充值", "转 USDT BEP20 并提交交易哈希。"),
        ("3. 购买配套", "选择 $500 到 $10,000 的投资配套。"),
        ("4. 选择业务", "把单位分配到 KTV、Beauty、Game House、Live House 或 Pet Cafe。"),
        ("5. 每日领取", "每天在后台领取项目 Token。"),
        ("6. 收益 + 管理", "查看 Token 价值、推荐奖励、40% 利润池、QR 消费与提现。"),
    ]
    for i, (h, b) in enumerate(steps):
        card(slide, 0.85, 1.32 + i * 0.78, 10.9, 0.58, h, b, [CYAN, GOLD, GREEN, PINK, VIOLET, CYAN][i], 8)
    tx(slide, DISCLAIMER, 0.85, 6.18, 10.8, 0.45, 8, MUTED)
    footer(slide)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    cover(prs)
    package_overview(prs)
    choose_business(prs)
    token_engine(prs)
    profit_pool(prs)
    referral_system(prs)
    topup_flow(prs)
    purchase_flow(prs)
    dashboard_usage(prs)
    risk_controls(prs)
    next_steps(prs)
    OUT.mkdir(exist_ok=True)
    output = OUT / "Reborn-Wave-Normal-Investor-Guide-ZH.pptx"
    prs.save(output)
    return output


if __name__ == "__main__":
    print(build())
