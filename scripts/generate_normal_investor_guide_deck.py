from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "investor-decks"
ASSETS = ROOT / "attached_assets"

BG = RGBColor(3, 21, 26)
PANEL = RGBColor(10, 42, 48)
PANEL_2 = RGBColor(13, 58, 65)
GOLD = RGBColor(248, 212, 119)
CYAN = RGBColor(34, 211, 238)
GREEN = RGBColor(52, 211, 153)
PINK = RGBColor(244, 114, 182)
VIOLET = RGBColor(167, 139, 250)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(185, 202, 207)
LINE = RGBColor(31, 82, 90)
FONT = "Aptos"

DISCLAIMER = (
    "This deck explains the planned investor system. Returns, token value, referral rewards, and profit "
    "sharing are projections and are not guaranteed. Legal, tax, accounting, and licensing review is required."
)

PACKAGES = [
    ("$500", "$500 credit", "1 share", "10% referral", "-"),
    ("$1,000", "$1,000 credit", "2 shares", "10% referral", "2D1N company tour"),
    ("$3,000", "$3,000 credit", "6 shares", "10% referral", "3D2N visit + entertainment"),
    ("$5,000", "$7,000 credit", "10 shares", "10% referral", "3D2N visit + entertainment"),
    ("$10,000", "$15,000 credit", "20 shares", "10% referral", "3D2N visit + entertainment"),
]

PROJECTS = [
    ("KTV Lounge", "Singing, rooms, hosting, competitions"),
    ("Beauty", "Salon services and glow-up packages"),
    ("Game House", "Kids, family games, redemption play"),
    ("Live House", "Sea-view band, dance, performances"),
    ("Pet Cafe", "Food, drinks, pets, blindbox content"),
]


def tx(slide, text, x, y, w, h, size=18, color=WHITE, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def rect(slide, x, y, w, h, fill=PANEL, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line or LINE
    shp.line.transparency = 22
    return shp


def image(slide, filename, x, y, w, h=None):
    path = ASSETS / filename
    if not path.exists():
        return
    if h:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    rect(slide, 10.55, -0.2, 2.4, 7.9, RGBColor(8, 55, 61), RGBColor(8, 55, 61), False)
    for i, color in enumerate([CYAN, GOLD, GREEN]):
        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2 + i * 0.55), Inches(0.45 + i), Inches(2.3), Inches(2.3))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.fill.transparency = 83
        shp.line.transparency = 100


def footer(slide):
    tx(slide, "Reborn Wave Group | Investor Package + Token System", 0.55, 7.05, 7.1, 0.2, 7, MUTED)
    tx(slide, "Projected returns only, not guaranteed.", 8.3, 6.95, 4.3, 0.24, 6, MUTED, align=PP_ALIGN.RIGHT)


def title(slide, heading, sub=None):
    tx(slide, "REBORN WAVE GROUP | NORMAL INVESTOR GUIDE", 0.55, 0.25, 4.8, 0.28, 9, MUTED, True)
    tx(slide, heading, 0.55, 0.62, 9.4, 0.72, 30, WHITE, True)
    if sub:
        tx(slide, sub, 0.58, 1.28, 9.2, 0.5, 12, MUTED)


def card(slide, x, y, w, h, heading, body, accent=CYAN, body_size=9):
    rect(slide, x, y, w, h, PANEL)
    tx(slide, heading, x + 0.18, y + 0.14, w - 0.36, 0.32, 13, accent, True)
    tx(slide, body, x + 0.18, y + 0.54, w - 0.36, h - 0.62, body_size, MUTED)


def metric(slide, x, y, w, label, value, accent=GOLD):
    rect(slide, x, y, w, 0.78, RGBColor(7, 39, 45), accent)
    tx(slide, value, x + 0.12, y + 0.12, w - 0.24, 0.28, 16, accent, True, PP_ALIGN.CENTER)
    tx(slide, label, x + 0.12, y + 0.45, w - 0.24, 0.2, 7, MUTED, False, PP_ALIGN.CENTER)


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    image(slide, "rwg-logo.png", 0.65, 0.42, 1.18)
    tx(slide, "Reborn Wave Group", 0.72, 1.55, 6.2, 0.7, 38, WHITE, True)
    tx(slide, "$500 - $10,000 Investor Package Guide", 0.76, 2.44, 6.5, 0.45, 18, GOLD, True)
    tx(slide, "Choose the business you believe in, receive daily project tokens, share in monthly business profit pools, and grow token value as more investors join.", 0.76, 3.14, 5.95, 1.05, 13, WHITE)
    rect(slide, 7.05, 0.9, 5.15, 5.6, PANEL_2)
    image(slide, "image_1764558901062.png", 7.38, 1.14, 4.45, 2.42)
    image(slide, "doluruu-blindbox-box.jpeg", 7.55, 3.88, 2.1, 1.22)
    image(slide, "doluruu-female-transparent.png", 10.0, 3.6, 1.12)
    tx(slide, "Portal: rebornwave.group/investor", 7.42, 5.55, 4.45, 0.35, 14, GOLD, True, PP_ALIGN.CENTER)
    footer(slide)


def package_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "Investor Packages", "Each package gives club spending credit, project shares, referral bonus, and selected travel/entertainment gifts.")
    headers = ["Amount", "Spending credit", "Shares", "Referral", "Extra gift"]
    widths = [1.35, 2.0, 1.35, 1.65, 4.15]
    x0, y0 = 0.82, 1.75
    x = x0
    for h, w in zip(headers, widths):
        rect(slide, x, y0, w, 0.5, RGBColor(8, 61, 68), CYAN, False)
        tx(slide, h, x + 0.05, y0 + 0.13, w - 0.1, 0.16, 8, WHITE, True, PP_ALIGN.CENTER)
        x += w
    for i, row in enumerate(PACKAGES):
        y = y0 + 0.55 + i * 0.62
        x = x0
        for c, (val, w) in enumerate(zip(row, widths)):
            fill = RGBColor(7, 39, 45) if i % 2 == 0 else RGBColor(9, 48, 54)
            rect(slide, x, y, w, 0.58, fill, LINE, False)
            color = GOLD if c in [1, 2] else WHITE if c == 0 else MUTED
            tx(slide, val, x + 0.06, y + 0.15, w - 0.12, 0.16, 8, color, c in [0, 1, 2], PP_ALIGN.CENTER if c < 3 else None)
            x += w
    card(slide, 0.82, 5.55, 10.65, 0.72, "Shared pool upside", "Each share is based on $500 participation for two years. 40% of the company shared pool is allocated to participating investors, with projected 5-10x upside potential, not guaranteed.", GOLD, 8)
    footer(slide)


def choose_business(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "Choose The Business To Invest In", "Investors vote with allocation units. The selected business decides which project tokens they earn daily.")
    for i, (name, desc) in enumerate(PROJECTS):
        card(slide, 0.72 + (i % 3) * 4.05, 1.75 + (i // 3) * 1.65, 3.55, 1.18, name, desc, [CYAN, PINK, GOLD, VIOLET, GREEN][i], 9)
    card(slide, 0.72, 5.18, 11.25, 0.86, "Allocation rule", "Investor package units must be fully allocated. Example: $10,000 Anchor Pack gives 20 units; investor chooses how many units go to each business.", GOLD, 9)
    footer(slide)


def token_engine(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "Daily Tokens + Growing Token Value", "Tokens reward participation and connect investors to the businesses they choose.")
    metric(slide, 0.85, 1.88, 2.4, "Daily token rule", "1 unit = 1 token/day", CYAN)
    metric(slide, 3.65, 1.88, 2.4, "Token price", "Pool / active tokens", GOLD)
    metric(slide, 6.45, 1.88, 2.4, "Package pool", "Value grows", GREEN)
    metric(slide, 9.25, 1.88, 2.4, "Sell option", "Token -> cash", PINK)
    card(slide, 0.85, 3.38, 3.35, 1.35, "Claim daily", "Investors press Claim Daily Project Tokens in the dashboard. Tokens match their active allocation units.", CYAN, 9)
    card(slide, 4.55, 3.38, 3.35, 1.35, "Value can increase", "When more people join, part of package activity enters the global token pool, supporting token price.", GOLD, 9)
    card(slide, 8.25, 3.38, 3.35, 1.35, "Sell tokens", "Investors can sell eligible active tokens into the cash wallet when pool liquidity is available.", GREEN, 9)
    card(slide, 0.85, 5.25, 10.75, 0.72, "Simple example", "$5,000 Founder Pack = 10 units. If 4 units are in KTV and 6 in Live House, daily claim earns 4 KTV tokens + 6 Live House tokens.", VIOLET, 9)
    footer(slide)


def profit_pool(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "Monthly Business Profit Pool", "The business profit model rewards investors who voted for each business.")
    metric(slide, 0.95, 1.86, 2.55, "Monthly verified profit", "Business report", CYAN)
    metric(slide, 3.9, 1.86, 2.55, "Investor pool", "40%", GOLD)
    metric(slide, 6.85, 1.86, 2.55, "Distribution", "To voters", GREEN)
    metric(slide, 9.8, 1.86, 1.7, "Basis", "Units", PINK)
    card(slide, 0.85, 3.35, 3.35, 1.4, "Who receives", "Investors who allocated units to that business are treated as voters for that business.", CYAN, 9)
    card(slide, 4.55, 3.35, 3.35, 1.4, "How it splits", "40% of monthly net profit is split by active units/votes for that business.", GOLD, 9)
    card(slide, 8.25, 3.35, 3.35, 1.4, "What admin checks", "Sales, operating costs, net profit, investor unit totals, and payout approval.", GREEN, 9)
    card(slide, 0.85, 5.25, 10.75, 0.72, "Example", "If Pet Cafe has $10,000 verified monthly net profit, $4,000 goes to the Pet Cafe voter pool and is split by Pet Cafe allocation units.", VIOLET, 9)
    footer(slide)


def referral_system(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "Referral System", "Investors can invite others and earn rewards when referred investors purchase packages.")
    metric(slide, 0.95, 1.86, 2.55, "Level 1 reward", "10%", GOLD)
    metric(slide, 3.9, 1.86, 2.55, "Level 2 reward", "5%", CYAN)
    metric(slide, 6.85, 1.86, 2.55, "Reward wallets", "Cash + spending", GREEN)
    metric(slide, 9.8, 1.86, 1.7, "Referral code", "INV-xxxx", PINK)
    card(slide, 0.85, 3.35, 3.35, 1.35, "Level 1", "Direct inviter receives referral reward when their invited investor buys a package.", GOLD, 9)
    card(slide, 4.55, 3.35, 3.35, 1.35, "Level 2", "Second-level inviter can receive an additional smaller reward when eligible.", CYAN, 9)
    card(slide, 8.25, 3.35, 3.35, 1.35, "Pool support", "Unassigned referral reward can flow back into the token pool, supporting token value.", GREEN, 9)
    card(slide, 0.85, 5.2, 10.75, 0.78, "Reward split", "The system records referral rewards in cash wallet and spending wallet so investors can withdraw part and use part inside the club.", VIOLET, 9)
    footer(slide)


def topup_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "How To Top Up", "The portal supports USDT BEP20 top-up and admin-confirmed crediting.")
    steps = [
        ("1. Register / login", "Open rebornwave.group/investor/login and enter the investor dashboard."),
        ("2. Choose package", "Pick $500, $1,000, $3,000, $5,000, or $10,000 package."),
        ("3. Connect wallet", "Connect BEP20 wallet or use manual payment proof if admin allows."),
        ("4. Send USDT", "Send the exact package amount to the Reborn Wave admin wallet."),
        ("5. Submit TX hash", "Paste transaction hash in dashboard for confirmation."),
        ("6. Cash wallet credited", "Once confirmed, cash balance is credited and can be used to purchase the package."),
    ]
    for i, (h, b) in enumerate(steps):
        card(slide, 0.72 + (i % 3) * 4.05, 1.75 + (i // 3) * 1.72, 3.55, 1.22, h, b, [CYAN, GOLD, GREEN, PINK, VIOLET, CYAN][i], 8)
    card(slide, 0.72, 5.48, 11.25, 0.78, "Safety note", "Only use the wallet address shown in the investor dashboard or confirmed by official Reborn Wave admin.", GOLD, 9)
    footer(slide)


def purchase_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "How To Buy A Package + Allocate Units", "After top-up, the investor purchases a package and chooses project allocation.")
    card(slide, 0.72, 1.78, 3.55, 1.32, "Select package", "Choose the package amount. The system calculates units automatically at $500 per unit.", CYAN, 9)
    card(slide, 4.55, 1.78, 3.55, 1.32, "Allocate units", "Choose how many units go into each business project before purchase confirmation.", GOLD, 9)
    card(slide, 8.38, 1.78, 3.55, 1.32, "Confirm purchase", "Cash wallet pays for the package and the dashboard records purchase + allocations.", GREEN, 9)
    card(slide, 0.72, 3.68, 5.45, 1.36, "Example A", "$1,000 Builder Pack = 2 units. Investor chooses 1 KTV unit and 1 Pet Cafe unit.", PINK, 10)
    card(slide, 6.48, 3.68, 5.45, 1.36, "Example B", "$10,000 Anchor Pack = 20 units. Investor chooses 8 Live House, 6 KTV, 4 Pet Cafe, 2 Beauty.", VIOLET, 10)
    card(slide, 0.72, 5.55, 11.25, 0.72, "After purchase", "Investors can claim daily project tokens, track token balances, view referrals, sell tokens, request withdrawals, and spend by QR wallet.", CYAN, 9)
    footer(slide)


def dashboard_usage(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "How To Use The Investor Dashboard", "The dashboard is where investors manage top-up, packages, tokens, spending, referrals, and withdrawals.")
    items = [
        ("Packages", "Choose package and allocate units to selected businesses."),
        ("Top-ups", "Credit wallet using USDT transaction hash or admin-confirmed payment."),
        ("Tokens", "Claim daily project tokens and see active/sold token balances."),
        ("Sell tokens", "Sell eligible tokens to cash wallet based on current pool price."),
        ("QR spending", "Use spending wallet inside the club with investor QR code."),
        ("Withdrawals", "Request cash wallet withdrawal to approved payout wallet."),
    ]
    for i, (h, b) in enumerate(items):
        card(slide, 0.72 + (i % 3) * 4.05, 1.75 + (i // 3) * 1.72, 3.55, 1.22, h, b, [CYAN, GOLD, GREEN, PINK, VIOLET, CYAN][i], 9)
    card(slide, 0.72, 5.5, 11.25, 0.72, "Dashboard link", "Investor home: rebornwave.group/investor   |   Login/Register: rebornwave.group/investor/login   |   Dashboard: rebornwave.group/investor/dashboard", GOLD, 9)
    footer(slide)


def risk_controls(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "Rules, Reports, And Risk Controls", "Clear records protect investors, operators, and Reborn Wave Group.")
    card(slide, 0.72, 1.78, 3.55, 1.35, "Monthly report", "Profit sharing should use verified revenue, costs, net profit, and approved payout calculations.", CYAN, 9)
    card(slide, 4.55, 1.78, 3.55, 1.35, "Token pool", "Token value depends on pool balance and active token supply. Price can rise or fall.", GOLD, 9)
    card(slide, 8.38, 1.78, 3.55, 1.35, "Admin approval", "Top-ups, withdrawals, package settings, and business reports need admin controls.", GREEN, 9)
    card(slide, 0.72, 3.68, 5.45, 1.36, "No guarantee", "More investors joining can increase pool value, but returns are not guaranteed and depend on actual system activity.", PINK, 9)
    card(slide, 6.48, 3.68, 5.45, 1.36, "Professional review", "Before public launch, contract, securities rules, tax, accounting, licenses, and investor risk disclosure must be reviewed.", VIOLET, 9)
    tx(slide, DISCLAIMER, 0.85, 5.72, 10.8, 0.48, 8, MUTED)
    footer(slide)


def next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, "Investor Journey Summary", "From new account to package purchase and daily rewards.")
    steps = [
        ("1. Create account", "Register at rebornwave.group/investor/login."),
        ("2. Top up", "Send USDT BEP20 and submit transaction hash."),
        ("3. Buy package", "Choose $500 to $10,000 package."),
        ("4. Pick businesses", "Allocate units to KTV, Beauty, Game House, Live House, or Pet Cafe."),
        ("5. Claim daily", "Collect project tokens every day from the dashboard."),
        ("6. Earn + manage", "Track token value, referrals, 40% profit pools, QR spending, and withdrawals."),
    ]
    for i, (h, b) in enumerate(steps):
        card(slide, 0.85, 1.32 + i * 0.78, 10.9, 0.58, h, b, [CYAN, GOLD, GREEN, PINK, VIOLET, CYAN][i], 8)
    tx(slide, DISCLAIMER, 0.85, 6.18, 10.8, 0.45, 8, MUTED)
    footer(slide)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    cover(prs)
    package_overview(prs)
    choose_business(prs)
    token_engine(prs)
    profit_pool(prs)
    referral_system(prs)
    topup_flow(prs)
    purchase_flow(prs)
    dashboard_usage(prs)
    risk_controls(prs)
    next_steps(prs)
    OUT.mkdir(exist_ok=True)
    output = OUT / "Reborn-Wave-Normal-Investor-Guide-EN.pptx"
    prs.save(output)
    return output


if __name__ == "__main__":
    print(build())
