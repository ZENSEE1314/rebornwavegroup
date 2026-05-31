from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "investor-decks"
ASSETS = ROOT / "attached_assets"
MAP_URL = (
    "https://www.google.com/maps/search/?api=1&query=Ruko%20Oceanic%20Bliss%2C%20Jl.%20Pasir%20Putih%20Harbourfront%20-%20Batam%20Centre.49%20blok%20A.%2051%2C%20Sadai%2C%20Bengkong%2C%20Batam%20City%2C%20Riau%20Islands%2029444"
)
ADDRESS = "Ruko Oceanic Bliss, Jl. Pasir Putih Harbourfront - Batam Centre.49 blok A. 51, Sadai, Bengkong, Batam City, Riau Islands 29444"

BG = RGBColor(3, 21, 26)
PANEL = RGBColor(10, 42, 48)
PANEL_2 = RGBColor(13, 58, 65)
GOLD = RGBColor(248, 212, 119)
CYAN = RGBColor(34, 211, 238)
GREEN = RGBColor(52, 211, 153)
PINK = RGBColor(244, 114, 182)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(185, 202, 207)


def tx(slide, text, x, y, w, h, size=24, color=WHITE, bold=False, font="Aptos", align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def rect(slide, x, y, w, h, fill=PANEL, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line or RGBColor(31, 82, 90)
    shp.line.transparency = 25
    return shp


def title(slide, deck, text, sub=None):
    tx(slide, deck["brand"], 0.55, 0.25, 3.0, 0.28, 9, MUTED, True, deck["font"])
    tx(slide, text, 0.55, 0.63, 8.7, 0.72, 31, WHITE, True, deck["font"])
    if sub:
        tx(slide, sub, 0.58, 1.32, 8.3, 0.5, 12, MUTED, False, deck["font"])


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    rect(slide, 10.55, -0.2, 2.4, 7.9, RGBColor(8, 55, 61), RGBColor(8, 55, 61), False).rotation = 0
    for i, color in enumerate([CYAN, GOLD, GREEN]):
        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2 + i * 0.55), Inches(0.45 + i * 1.0), Inches(2.3), Inches(2.3))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.fill.transparency = 83
        shp.line.transparency = 100


def image(slide, path, x, y, w, h=None):
    p = ASSETS / path
    if p.exists():
        if h:
            slide.shapes.add_picture(str(p), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
        else:
            slide.shapes.add_picture(str(p), Inches(x), Inches(y), width=Inches(w))


def footer(slide, deck):
    tx(slide, deck["footer"], 0.55, 7.0, 7.8, 0.2, 7, MUTED, False, deck["font"])
    tx(slide, deck["disclaimer_short"], 8.4, 6.88, 4.35, 0.34, 6, MUTED, False, deck["font"], PP_ALIGN.RIGHT)


def card(slide, deck, x, y, w, h, heading, body, accent=CYAN):
    rect(slide, x, y, w, h, PANEL)
    tx(slide, heading, x + 0.18, y + 0.15, w - 0.36, 0.33, 14, accent, True, deck["font"])
    tx(slide, body, x + 0.18, y + 0.56, w - 0.36, h - 0.72, 9, MUTED, False, deck["font"])


def cover(prs, deck):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    image(slide, "rwg-logo.png", 0.65, 0.45, 1.15)
    tx(slide, deck["cover_title"], 0.7, 1.75, 6.0, 1.35, 40, WHITE, True, deck["font"])
    tx(slide, deck["cover_sub"], 0.75, 3.12, 5.8, 0.8, 15, MUTED, False, deck["font"])
    tx(slide, deck["cover_tag"], 0.76, 4.18, 4.8, 0.32, 13, GOLD, True, deck["font"])
    rect(slide, 7.05, 0.9, 5.15, 5.6, PANEL_2)
    image(slide, "doluruu-blindbox-box.jpeg", 7.38, 1.22, 4.45, 2.55)
    image(slide, "doluruu-female-transparent.png", 8.25, 3.42, 1.45)
    image(slide, "Doluruu Boy_1749664545355.png", 9.62, 3.35, 1.32)
    tx(slide, deck["cover_box"], 7.45, 5.5, 4.35, 0.35, 15, GOLD, True, deck["font"], PP_ALIGN.CENTER)
    footer(slide, deck)


def opportunity(prs, deck):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, deck, deck["slides"][1][0], deck["slides"][1][1])
    items = deck["opportunity"]
    for i, item in enumerate(items):
        card(slide, deck, 0.68 + i * 3.05, 2.0, 2.75, 2.55, item[0], item[1], [CYAN, GOLD, GREEN][i])
    rect(slide, 10.0, 1.8, 2.35, 3.0, RGBColor(8, 61, 68))
    tx(slide, "$200K", 10.22, 2.18, 1.9, 0.55, 28, GOLD, True, deck["font"], PP_ALIGN.CENTER)
    tx(slide, deck["cap_label"], 10.15, 2.92, 2.0, 0.8, 12, WHITE, True, deck["font"], PP_ALIGN.CENTER)
    tx(slide, deck["cap_note"], 10.15, 3.78, 2.0, 0.55, 8, MUTED, False, deck["font"], PP_ALIGN.CENTER)
    footer(slide, deck)


def floors(prs, deck):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, deck, deck["slides"][2][0], deck["slides"][2][1])
    colors = [CYAN, PINK, GOLD, GREEN, RGBColor(167, 139, 250)]
    for i, fl in enumerate(deck["floors"]):
        y = 1.75 + i * 0.86
        rect(slide, 0.75, y, 11.6, 0.63, RGBColor(7, 39, 45), colors[i])
        tx(slide, fl[0], 0.95, y + 0.13, 0.85, 0.23, 15, colors[i], True, deck["font"])
        tx(slide, fl[1], 1.78, y + 0.1, 3.4, 0.3, 13, WHITE, True, deck["font"])
        tx(slide, fl[2], 5.25, y + 0.12, 6.7, 0.25, 9, MUTED, False, deck["font"])
    footer(slide, deck)


def investor(prs, deck):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, deck, deck["slides"][3][0], deck["slides"][3][1])
    for i, item in enumerate(deck["investor_steps"]):
        card(slide, deck, 0.65 + (i % 2) * 5.95, 1.8 + (i // 2) * 1.75, 5.55, 1.35, item[0], item[1], [GOLD, CYAN, GREEN, PINK][i])
    tx(slide, deck["investor_warning"], 0.8, 5.7, 11.2, 0.6, 10, MUTED, False, deck["font"])
    footer(slide, deck)


def referral(prs, deck):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, deck, deck["slides"][4][0], deck["slides"][4][1])
    for i, item in enumerate(deck["referral"]):
        x = 0.75 + i * 2.35
        rect(slide, x, 2.1, 2.0, 2.35, RGBColor(7, 39, 45), [CYAN, GOLD, GREEN, PINK, CYAN][i])
        tx(slide, item[0], x + 0.2, 2.48, 1.6, 0.42, 24, [CYAN, GOLD, GREEN, PINK, CYAN][i], True, deck["font"], PP_ALIGN.CENTER)
        tx(slide, item[1], x + 0.18, 3.25, 1.64, 0.7, 9, MUTED, False, deck["font"], PP_ALIGN.CENTER)
    tx(slide, deck["referral_note"], 0.82, 5.35, 10.8, 0.7, 12, WHITE, True, deck["font"], PP_ALIGN.CENTER)
    footer(slide, deck)


def blindbox(prs, deck):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, deck, deck["slides"][5][0], deck["slides"][5][1])
    image(slide, "Doluruu Boy_1749664545355.png", 0.85, 2.1, 1.25)
    image(slide, "doluruu-female-transparent.png", 2.35, 2.05, 1.28)
    image(slide, "Doluruu Baby_1749663725243.png", 3.82, 2.22, 1.05)
    image(slide, "doluruu-blindbox-box.jpeg", 5.35, 2.0, 2.65, 1.65)
    for i, item in enumerate(deck["blindbox_rules"]):
        card(slide, deck, 8.35, 1.55 + i * 1.05, 3.65, 0.8, item[0], item[1], [GOLD, CYAN, GREEN, PINK][i])
    tx(slide, deck["token_line"], 0.95, 4.8, 6.8, 0.55, 17, GOLD, True, deck["font"])
    footer(slide, deck)


def marketing(prs, deck):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, deck, deck["slides"][6][0], deck["slides"][6][1])
    for i, item in enumerate(deck["agents"]):
        card(slide, deck, 0.65 + (i % 3) * 4.05, 1.72 + (i // 3) * 1.65, 3.65, 1.22, item[0], item[1], [GOLD, CYAN, GREEN, PINK, RGBColor(167, 139, 250), CYAN][i])
    footer(slide, deck)


def location(prs, deck):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, deck, deck["slides"][7][0], deck["slides"][7][1])
    rect(slide, 0.75, 1.85, 6.1, 2.65, PANEL)
    tx(slide, "Ruko Oceanic Bliss", 1.05, 2.18, 5.3, 0.38, 20, GOLD, True, deck["font"])
    tx(slide, ADDRESS, 1.05, 2.85, 5.3, 0.92, 13, WHITE, False, deck["font"])
    link = tx(slide, deck["map_cta"], 1.05, 3.95, 3.5, 0.35, 14, CYAN, True, deck["font"])
    link.click_action.hyperlink.address = MAP_URL
    rect(slide, 7.32, 1.85, 4.55, 2.65, RGBColor(8, 61, 68))
    tx(slide, deck["map_symbol"], 8.25, 2.35, 2.8, 0.8, 40, GOLD, True, deck["font"], PP_ALIGN.CENTER)
    tx(slide, deck["map_cta"], 8.0, 3.35, 3.25, 0.35, 15, WHITE, True, deck["font"], PP_ALIGN.CENTER).click_action.hyperlink.address = MAP_URL
    footer(slide, deck)


def next_steps(prs, deck):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    title(slide, deck, deck["slides"][8][0], deck["slides"][8][1])
    for i, item in enumerate(deck["next_steps"]):
        card(slide, deck, 0.75, 1.75 + i * 1.03, 11.25, 0.75, item[0], item[1], [CYAN, GOLD, GREEN, PINK][i % 4])
    tx(slide, deck["disclaimer"], 0.85, 6.03, 10.9, 0.58, 8, MUTED, False, deck["font"])
    footer(slide, deck)


def build(deck, filename):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    cover(prs, deck)
    opportunity(prs, deck)
    floors(prs, deck)
    investor(prs, deck)
    referral(prs, deck)
    blindbox(prs, deck)
    marketing(prs, deck)
    location(prs, deck)
    next_steps(prs, deck)
    OUT.mkdir(exist_ok=True)
    prs.save(OUT / filename)


EN = {
    "font": "Aptos",
    "brand": "REBORN WAVE GROUP | INVESTOR DECK",
    "footer": "KTV | Game House | Beauty | Pet Cafe | Sea-View Live House | Blindbox Rewards",
    "disclaimer_short": "ROI target, not guaranteed return.",
    "cover_title": "Reborn Wave Group Club Economy",
    "cover_sub": "A 5-story waterfront lifestyle ecosystem in Batam built for customers, members, investors, partners, and workers.",
    "cover_tag": "$5,000 investor package | $200K capped allocation",
    "cover_box": "Every investor receives 1 free Doluruu blindbox",
    "slides": [
        (),
        ("The Opportunity", "One venue with multiple spending reasons across day, family, beauty, community, nightlife, and investor-hosted visits."),
        ("Five Floors, One Customer Loop", "Each floor gives people another reason to visit, stay longer, post content, and return."),
        ("Investor Package", "$5,000 enters the ecosystem as investor capital and customer-hosting credits."),
        ("Referral Logic", "Referral bonuses reward growth and reduce the remaining ROI target."),
        ("Blindbox Pet Economy", "Doluruu pets create daily token habits and prize exchange reasons."),
        ("Marketing Agent Team", "Ruflo-style agents coordinate content, engagement, investor leads, partnerships, and analytics."),
        ("Location", "A clear Batam destination with map-click convenience for visitors and investor hosting."),
        ("Next Steps", "Move from interest to launch-ready execution with legal review, campaign assets, and partner outreach."),
    ],
    "opportunity": [
        ("Customer Traffic", "KTV, kids games, beauty, pet cafe, live music, dance floor, and sea view create multiple audiences."),
        ("Investor Flywheel", "Credits help investors bring friends, customers, and partners into the club experience."),
        ("Member Retention", "Blindbox pet feeding, daily tokens, prizes, competitions, and events keep people returning."),
    ],
    "cap_label": "Total investor allocation cap",
    "cap_note": "Scarcity supports urgency and clean package communication.",
    "floors": [
        ("1F", "KTV Lounge + Kids Game House", "2 units KTV lounge, 1 unit kids game house for singing competitions and family traffic."),
        ("2F", "Private KTV + Beauty Salon", "4 private KTV rooms plus beauty services for repeat bookings."),
        ("3F", "VIP KTV + Beauty Salon", "2 VIP rooms for investor guests, birthdays, business hosting, and premium spend."),
        ("4F", "Pet Cafe", "All 3 units become a pet cafe tied to Doluruu, content, food, drinks, and community."),
        ("5F", "Sea-View Live House", "Live band, dance floor, rooftop sea view, event nights, and premium table sales."),
    ],
    "investor_steps": [
        ("Invest $5,000", "Investor receives $5,000 club spending credits for friends, customers, and private hosting."),
        ("Target $10,000 ROI", "Target return across 2 years under contract terms and monthly payout rules."),
        ("20% Profit Pool", "20% of business profit is allocated monthly and divided among active investors."),
        ("10% New Investor Pool", "10% of every new $5,000 package is shared with earlier active investors."),
    ],
    "investor_warning": "Investor terms require qualified legal and financial review before public launch or contract signing.",
    "referral": [
        ("1", "$1,000 referral bonus paid; ROI target reduced by $1,000."),
        ("3", "$3,000 bonus total; remaining target becomes $7,000."),
        ("5", "$5,000 bonus total; remaining target becomes $5,000."),
        ("10", "$10,000 bonus total; ROI target cleared."),
        ("Again", "Investor may join again after contract completion."),
    ],
    "referral_note": "Each invited friend invests $5,000. The inviter receives $1,000, deducted from the inviter's remaining ROI target.",
    "blindbox_rules": [
        ("Investor Gift", "1 free blindbox per investor package."),
        ("Male + Female", "Owning both unlocks 1 baby pet free."),
        ("Daily Tokens", "Feed each pet daily to earn tokens."),
        ("Prize Exchange", "Tokens can support rewards, upgrades, and club campaigns."),
    ],
    "token_line": "3 pets = 3 daily tokens: male + female + baby.",
    "agents": [
        ("Growth Queen", "Sets weekly priorities, assigns agents, and checks compliance."),
        ("Content Producer", "Creates Reels, carousels, captions, and multilingual posts."),
        ("Social Amplifier", "Runs polls, comments, UGC prompts, and share mechanics."),
        ("Investor Leads", "Builds investor posts, WhatsApp scripts, FAQs, and follow-ups."),
        ("Partnerships", "Targets creators, singers, salons, hotels, pet groups, and events."),
        ("Analytics Optimizer", "Tracks reach, shares, map clicks, inquiries, calls, and signups."),
    ],
    "map_cta": "Open Google Maps",
    "map_symbol": "BATAM",
    "next_steps": [
        ("1. Finalize contract", "Legal review for ROI target language, payout rules, risk disclosure, and referral deductions."),
        ("2. Launch content engine", "Daily short-form video, floor tours, Doluruu posts, and event content."),
        ("3. Open investor conversations", "Use the deck, website, WhatsApp script, and showroom hosting credits."),
        ("4. Recruit the floor team", "Hosts, singers, stylists, pet cafe crew, event staff, and promoters."),
    ],
    "disclaimer": "This deck is a business and marketing draft. ROI references are targets, not guaranteed returns. Investment contracts, payout rules, profit calculations, referral deductions, compliance requirements, and risk disclosures should be reviewed by qualified legal and financial professionals before launch.",
}


ZH = {
    **EN,
    "font": "Microsoft YaHei",
    "brand": "REBORN WAVE GROUP | 投资者简报",
    "footer": "KTV | 游戏屋 | 美容 | 宠物咖啡厅 | 海景 Live House | 盲盒奖励",
    "disclaimer_short": "ROI 为目标，并非保证收益。",
    "cover_title": "Reborn Wave Group 俱乐部经济生态",
    "cover_sub": "位于巴淡岛海边的五层生活娱乐生态，连接顾客、会员、投资者、合作伙伴和员工。",
    "cover_tag": "$5,000 投资配套 | 总额度上限 $200K",
    "cover_box": "每位投资者免费获得 1 个 Doluruu 盲盒",
    "slides": [
        (),
        ("市场机会", "一个场地覆盖白天家庭、美容、社群、夜生活和投资者招待等多种消费理由。"),
        ("五层空间，一个消费循环", "每一层都给顾客更多到店、停留、拍内容和复访的理由。"),
        ("投资者配套", "$5,000 进入生态，同时形成投资资金和顾客招待消费额度。"),
        ("推荐逻辑", "推荐奖金奖励增长，并扣减剩余 ROI 目标。"),
        ("盲盒宠物经济", "Doluruu 宠物带来每日喂养、代币和奖品兑换习惯。"),
        ("营销代理团队", "Ruflo 风格代理团队协调内容、互动、投资者线索、合作和数据分析。"),
        ("位置", "巴淡岛清晰到店地址，方便顾客和投资者招待。"),
        ("下一步", "通过法律审核、内容资产和合作拓展，把兴趣转化为启动执行。"),
    ],
    "opportunity": [
        ("顾客流量", "KTV、儿童游戏、美容、宠物咖啡、现场音乐、舞池和海景覆盖多种客群。"),
        ("投资飞轮", "消费额度帮助投资者带朋友、客户和伙伴进入俱乐部体验。"),
        ("会员留存", "盲盒宠物喂养、每日代币、奖品、比赛和活动让会员持续回来。"),
    ],
    "cap_label": "总投资额度上限",
    "cap_note": "稀缺额度让配套更清晰，也支持行动紧迫感。",
    "floors": [
        ("1F", "KTV 大厅 + 儿童游戏屋", "2 个单位 KTV 大厅，1 个单位儿童游戏屋，支持唱歌比赛和家庭客流。"),
        ("2F", "私人 KTV + 美容院", "4 间私人 KTV 加美容服务，带动重复预订。"),
        ("3F", "VIP KTV + 美容院", "2 间 VIP 房，适合投资者招待、生日、商务聚会和高消费。"),
        ("4F", "宠物咖啡厅", "三个单位打通，连接 Doluruu、内容、餐饮和宠物社群。"),
        ("5F", "海景 Live House", "现场乐队、舞池、海景、活动夜和高端桌位销售。"),
    ],
    "investor_steps": [
        ("投资 $5,000", "投资者获得 $5,000 俱乐部消费额度，可用于朋友、客户和私人招待。"),
        ("目标 $10,000 ROI", "目标回报在 2 年内按正式合同和月度派发规则执行。"),
        ("20% 利润池", "每月将业务利润的 20% 分配给活跃投资者。"),
        ("10% 新投资池", "每个新的 $5,000 投资配套中，10% 分享给较早加入的活跃投资者。"),
    ],
    "investor_warning": "投资条款在公开发布或签署合同前，需要合资格法律和财务专业人士审核。",
    "referral": [
        ("1 位", "支付 $1,000 推荐奖金，ROI 目标扣减 $1,000。"),
        ("3 位", "奖金累计 $3,000，剩余目标 $7,000。"),
        ("5 位", "奖金累计 $5,000，剩余目标 $5,000。"),
        ("10 位", "奖金累计 $10,000，ROI 目标清零。"),
        ("再加入", "合约完成后可再次投资。"),
    ],
    "referral_note": "每位被邀请朋友投资 $5,000，邀请人获得 $1,000，该金额从邀请人的剩余 ROI 目标中扣减。",
    "blindbox_rules": [
        ("投资者礼物", "每个投资配套免费获得 1 个盲盒。"),
        ("公 + 母", "同时拥有公母宠物，可免费获得 1 个宝宝。"),
        ("每日代币", "每天喂养宠物即可获得代币。"),
        ("奖品兑换", "代币可用于奖励、升级和俱乐部活动。"),
    ],
    "token_line": "3 只宠物 = 每天 3 个代币：公 + 母 + 宝宝。",
    "agents": [
        ("增长 Queen", "设定每周优先级、分配代理并检查合规。"),
        ("内容制作", "制作短视频、轮播、文案和多语言帖子。"),
        ("社交放大", "设计投票、评论、UGC 和分享机制。"),
        ("投资线索", "制作投资帖、WhatsApp 话术、FAQ 和跟进流程。"),
        ("合作拓展", "寻找创作者、歌手、美容、酒店、宠物社群和活动伙伴。"),
        ("数据优化", "追踪曝光、分享、地图点击、询盘、通话和注册。"),
    ],
    "map_cta": "打开 Google 地图",
    "map_symbol": "巴淡岛",
    "next_steps": [
        ("1. 完成合约", "法律审核 ROI 目标语言、派发规则、风险披露和推荐扣减。"),
        ("2. 启动内容引擎", "每日短视频、楼层导览、Doluruu 内容和活动内容。"),
        ("3. 开始投资者沟通", "使用简报、网站、WhatsApp 话术和到店招待额度。"),
        ("4. 招募楼层团队", "主持、歌手、美容师、宠物咖啡员工、活动人员和推广员。"),
    ],
    "disclaimer": "本简报为业务和营销草稿。ROI 为目标，并非保证收益。投资合同、派发规则、利润计算、推荐扣减、合规要求和风险披露，应在上线前由合资格法律和财务专业人士审核。",
}


if __name__ == "__main__":
    build(EN, "Reborn-Wave-Investor-Deck-EN.pptx")
    build(ZH, "Reborn-Wave-Investor-Deck-ZH.pptx")
    print(f"Created decks in {OUT}")
